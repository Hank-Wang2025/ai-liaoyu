#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
更新PPT，导入毛泽东故居图片并排版对齐
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os


def update_ppt(ppt_path: str, image_path: str, output_path: str):
    """更新PPT，导入图片并排版对齐"""
    
    prs = Presentation(ppt_path)
    
    # 获取图片尺寸
    from PIL import Image
    with Image.open(image_path) as img:
        img_width, img_height = img.size
    
    # 计算图片比例
    img_ratio = img_width / img_height
    
    # ========== 第2页：故居介绍 - 替换左侧占位图 ==========
    slide = prs.slides[1]  # 第2页
    
    # 找到并删除原来的占位形状
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
            # 检查是否是占位图（通过位置和大小判断）
            if hasattr(shape, 'left') and shape.left < Inches(6):
                shapes_to_remove.append(shape)
    
    # 删除占位形状
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
    
    # 添加真实图片 - 左侧区域，对齐到指定位置
    # 原占位图位置：Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.5)
    left = Inches(0.8)
    top = Inches(1.5)
    max_width = Inches(5.5)
    max_height = Inches(5.5)
    
    # 根据图片比例计算实际尺寸
    if img_ratio > 1:  # 宽图
        pic_width = max_width
        pic_height = pic_width / img_ratio
        if pic_height > max_height:
            pic_height = max_height
            pic_width = pic_height * img_ratio
    else:  # 高图
        pic_height = max_height
        pic_width = pic_height * img_ratio
        if pic_width > max_width:
            pic_width = max_width
            pic_height = pic_width / img_ratio
    
    # 居中放置
    pic_left = left + (max_width - pic_width) / 2
    pic_top = top + (max_height - pic_height) / 2
    
    slide.shapes.add_picture(image_path, pic_left, pic_top, pic_width, pic_height)
    
    # ========== 第5页：作品展示 - 替换作品展示区 ==========
    slide = prs.slides[4]  # 第5页
    
    # 找到并删除原来的占位形状
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
            shapes_to_remove.append(shape)
    
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
    
    # 添加图片 - 作品展示区，居中显示
    # 原占位图位置：Inches(2), Inches(1.4), Inches(9.333), Inches(5.5)
    left = Inches(2)
    top = Inches(1.4)
    max_width = Inches(9.333)
    max_height = Inches(5.5)
    
    # 根据图片比例计算实际尺寸
    if img_ratio > max_width / max_height:  # 更宽的图
        pic_width = max_width
        pic_height = pic_width / img_ratio
    else:  # 更高的图
        pic_height = max_height
        pic_width = pic_height * img_ratio
    
    # 居中放置
    pic_left = left + (max_width - pic_width) / 2
    pic_top = top + (max_height - pic_height) / 2
    
    slide.shapes.add_picture(image_path, pic_left, pic_top, pic_width, pic_height)
    
    # 更新作品名称
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "作品名称" in text:
                shape.text_frame.text = "作品名称：《毛泽东故居》—— 手绘作品"
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.size = Pt(22)
                    paragraph.font.color.rgb = RGBColor(100, 100, 100)
                    paragraph.alignment = PP_ALIGN.CENTER
    
    # 保存
    prs.save(output_path)
    print(f"PPT已更新并保存至: {output_path}")


if __name__ == "__main__":
    desktop = os.path.expanduser("~/Desktop")
    ppt_path = os.path.join(desktop, "毛泽东故居-小学生美术作品.pptx")
    image_path = os.path.join(desktop, "mao_artwork_demo.jpg")
    output_path = os.path.join(desktop, "毛泽东故居-小学生美术作品(含图片).pptx")
    
    update_ppt(ppt_path, image_path, output_path)
