#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
创建小学生美术作品PPT - 毛泽东故居主题
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os


def create_ppt(output_path: str):
    """创建小学生美术作品PPT"""
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ========== 第1页：封面 ==========
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景色 - 温暖的米黄色
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 248, 230)
    background.line.fill.background()
    
    # 标题装饰框
    title_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5),
        Inches(10.333), Inches(2)
    )
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = RGBColor(200, 50, 50)
    title_box.line.fill.background()
    
    # 标题文字
    title_tf = title_box.text_frame
    title_tf.text = "走进毛泽东故居"
    title_tf.paragraphs[0].font.size = Pt(54)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle = slide.shapes.add_textbox(
        Inches(2), Inches(4), Inches(9.333), Inches(1)
    )
    subtitle_tf = subtitle.text_frame
    subtitle_tf.text = "——小学生美术作品展"
    subtitle_tf.paragraphs[0].font.size = Pt(36)
    subtitle_tf.paragraphs[0].font.color.rgb = RGBColor(139, 69, 19)
    subtitle_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 作者信息
    author = slide.shapes.add_textbox(
        Inches(2), Inches(5.5), Inches(9.333), Inches(1)
    )
    author_tf = author.text_frame
    author_tf.text = "班级：_______  \n姓名：_______  \n指导教师：_______"
    author_tf.paragraphs[0].font.size = Pt(24)
    author_tf.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
    author_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ========== 第2页：故居介绍 ==========
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(245, 245, 240)
    background.line.fill.background()
    
    # 页面标题
    header = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
    )
    header_tf = header.text_frame
    header_tf.text = "毛泽东故居简介"
    header_tf.paragraphs[0].font.size = Pt(40)
    header_tf.paragraphs[0].font.bold = True
    header_tf.paragraphs[0].font.color.rgb = RGBColor(180, 40, 40)
    
    # 分隔线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1),
        Inches(12.333), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(180, 40, 40)
    line.line.fill.background()
    
    # 内容区域 - 左侧图片占位
    img_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5),
        Inches(5.5), Inches(5.5)
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = RGBColor(230, 230, 220)
    img_placeholder.line.color.rgb = RGBColor(180, 40, 40)
    img_placeholder.line.width = Pt(3)
    
    # 图片占位文字
    img_text = img_placeholder.text_frame
    img_text.text = "【毛泽东故居照片】\n\n（可以粘贴参观时拍摄的照片）"
    img_text.paragraphs[0].font.size = Pt(20)
    img_text.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
    img_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 内容区域 - 右侧文字
    content = slide.shapes.add_textbox(
        Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5)
    )
    content_tf = content.text_frame
    content_tf.word_wrap = True
    
    intro_text = """📍 位置：湖南省湘潭市韶山市韶山冲

🏠 概况：
毛泽东故居是一座土木结构的"凹"字型建筑，坐南朝北，建筑面积约472平方米。

📅 历史：
1893年12月26日，毛泽东诞生在这里，并在此度过了童年和少年时代。

🎨 绘画要点：
• 土黄色的墙壁
• 青瓦屋顶
• 门前的池塘
• 周围的青山绿树"""
    
    content_tf.text = intro_text
    for paragraph in content_tf.paragraphs:
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = RGBColor(60, 60, 60)
        paragraph.space_after = Pt(12)
    
    # ========== 第3页：创作灵感 ==========
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 250, 245)
    background.line.fill.background()
    
    header = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
    )
    header_tf = header.text_frame
    header_tf.text = "我的创作灵感"
    header_tf.paragraphs[0].font.size = Pt(40)
    header_tf.paragraphs[0].font.bold = True
    header_tf.paragraphs[0].font.color.rgb = RGBColor(180, 40, 40)
    
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1),
        Inches(12.333), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(180, 40, 40)
    line.line.fill.background()
    
    # 三个灵感卡片
    card_colors = [RGBColor(255, 220, 200), RGBColor(220, 240, 220), RGBColor(220, 230, 255)]
    card_titles = ["🎨 色彩印象", "🏛️ 建筑特色", "🌳 自然环境"]
    card_contents = [
        "故居的土黄色墙壁给我留下深刻印象，\n它诉说着历史的沧桑。\n我决定用暖色调来表现\n这种古朴的美感。",
        "凹字形的建筑布局很有特色，\n青瓦白墙搭配和谐。\n我要画出这种独特的\n湖南民居风格。",
        "故居前有池塘，后有青山，\n环境优美。\n我想表现人与自然\n和谐相处的美好画面。"
    ]
    
    for i, (color, title, text) in enumerate(zip(card_colors, card_titles, card_contents)):
        x_pos = Inches(0.8 + i * 4.2)
        
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(1.5),
            Inches(3.8), Inches(5.3)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.color.rgb = RGBColor(180, 40, 40)
        card.line.width = Pt(2)
        
        card_tf = card.text_frame
        card_tf.text = f"{title}\n\n{text}"
        card_tf.paragraphs[0].font.size = Pt(24)
        card_tf.paragraphs[0].font.bold = True
        card_tf.paragraphs[0].font.color.rgb = RGBColor(180, 40, 40)
        card_tf.paragraphs[2].font.size = Pt(20)
        card_tf.paragraphs[2].font.color.rgb = RGBColor(80, 80, 80)
    
    # ========== 第4页：绘画过程 ==========
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(250, 250, 245)
    background.line.fill.background()
    
    header = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
    )
    header_tf = header.text_frame
    header_tf.text = "绘画过程"
    header_tf.paragraphs[0].font.size = Pt(40)
    header_tf.paragraphs[0].font.bold = True
    header_tf.paragraphs[0].font.color.rgb = RGBColor(180, 40, 40)
    
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1),
        Inches(12.333), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(180, 40, 40)
    line.line.fill.background()
    
    # 四个步骤
    steps = [
        ("第一步：铅笔起稿", "用铅笔轻轻勾勒出\n故居的轮廓和\n主要结构。"),
        ("第二步：上色", "先涂大面积的背景色，\n如天空、地面、\n墙面等。"),
        ("第三步：细节刻画", "画出瓦片、门窗、\n树木等细节，\n增加层次感。"),
        ("第四步：调整完善", "检查整体效果，\n补充阴影和\n高光部分。")
    ]
    
    for i, (title, desc) in enumerate(steps):
        row = i // 2
        col = i % 2
        x_pos = Inches(0.8 + col * 6.2)
        y_pos = Inches(1.4 + row * 3)
        
        # 步骤框
        step_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos,
            Inches(5.8), Inches(2.7)
        )
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        step_box.line.color.rgb = RGBColor(200, 150, 100)
        step_box.line.width = Pt(2)
        
        # 步骤编号圆圈
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x_pos + Inches(0.2), y_pos + Inches(0.5),
            Inches(0.8), Inches(0.8)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(180, 40, 40)
        circle.line.fill.background()
        
        circle_tf = circle.text_frame
        circle_tf.text = str(i + 1)
        circle_tf.paragraphs[0].font.size = Pt(28)
        circle_tf.paragraphs[0].font.bold = True
        circle_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        circle_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 步骤文字
        step_text = slide.shapes.add_textbox(
            x_pos + Inches(1.2), y_pos + Inches(0.3),
            Inches(4.4), Inches(2)
        )
        step_tf = step_text.text_frame
        step_tf.text = f"{title}\n{desc}"
        step_tf.paragraphs[0].font.size = Pt(22)
        step_tf.paragraphs[0].font.bold = True
        step_tf.paragraphs[0].font.color.rgb = RGBColor(180, 40, 40)
        step_tf.paragraphs[1].font.size = Pt(18)
        step_tf.paragraphs[1].font.color.rgb = RGBColor(80, 80, 80)
    
    # ========== 第5页：作品展示 ==========
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 248, 240)
    background.line.fill.background()
    
    header = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
    )
    header_tf = header.text_frame
    header_tf.text = "我的美术作品"
    header_tf.paragraphs[0].font.size = Pt(40)
    header_tf.paragraphs[0].font.bold = True
    header_tf.paragraphs[0].font.color.rgb = RGBColor(180, 40, 40)
    
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1),
        Inches(12.333), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(180, 40, 40)
    line.line.fill.background()
    
    # 作品展示区
    artwork_frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(1.4),
        Inches(9.333), Inches(5.5)
    )
    artwork_frame.fill.solid()
    artwork_frame.fill.fore_color.rgb = RGBColor(255, 255, 255)
    artwork_frame.line.color.rgb = RGBColor(180, 40, 40)
    artwork_frame.line.width = Pt(4)
    
    artwork_tf = artwork_frame.text_frame
    artwork_tf.text = "【请在此粘贴您的美术作品】\n\n\n\n\n（可以是水彩画、蜡笔画、素描等任何形式）"
    artwork_tf.paragraphs[0].font.size = Pt(24)
    artwork_tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
    artwork_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 作品名称
    caption = slide.shapes.add_textbox(
        Inches(2), Inches(6.8), Inches(9.333), Inches(0.5)
    )
    caption_tf = caption.text_frame
    caption_tf.text = "作品名称：《毛泽东故居》"
    caption_tf.paragraphs[0].font.size = Pt(22)
    caption_tf.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
    caption_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ========== 第6页：心得体会 ==========
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(250, 245, 240)
    background.line.fill.background()
    
    header = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
    )
    header_tf = header.text_frame
    header_tf.text = "我的创作心得"
    header_tf.paragraphs[0].font.size = Pt(40)
    header_tf.paragraphs[0].font.bold = True
    header_tf.paragraphs[0].font.color.rgb = RGBColor(180, 40, 40)
    
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1),
        Inches(12.333), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(180, 40, 40)
    line.line.fill.background()
    
    # 心得内容框
    thought_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.5),
        Inches(11.333), Inches(5.5)
    )
    thought_box.fill.solid()
    thought_box.fill.fore_color.rgb = RGBColor(255, 252, 245)
    thought_box.line.color.rgb = RGBColor(200, 150, 100)
    thought_box.line.width = Pt(2)
    
    thought_tf = thought_box.text_frame
    thought_tf.word_wrap = True
    
    template_text = """通过这次创作，我学到了：

    🎨 绘画技巧方面：
    • 如何运用透视原理表现建筑的立体感
    • 暖色调和冷色调的搭配使用
    • 水彩画的干湿画法技巧

    📚 历史文化方面：
    • 了解了毛泽东爷爷的成长故事
    • 认识到保护历史建筑的重要性
    • 感受到革命先辈艰苦奋斗的精神

    💡 个人感悟：
    （请写下你自己的感受......）"""
    
    thought_tf.text = template_text
    for paragraph in thought_tf.paragraphs:
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = RGBColor(60, 60, 60)
        paragraph.space_after = Pt(8)
    
    # ========== 第7页：结束页 ==========
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(180, 40, 40)
    background.line.fill.background()
    
    # 感谢文字
    thanks = slide.shapes.add_textbox(
        Inches(2), Inches(2.5), Inches(9.333), Inches(1.5)
    )
    thanks_tf = thanks.text_frame
    thanks_tf.text = "感谢观看"
    thanks_tf.paragraphs[0].font.size = Pt(60)
    thanks_tf.paragraphs[0].font.bold = True
    thanks_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    thanks_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 装饰图案 - 五角星
    for i in range(5):
        x_pos = Inches(2 + i * 2)
        star = slide.shapes.add_shape(
            MSO_SHAPE.STAR_5_POINT, x_pos, Inches(4.5),
            Inches(0.8), Inches(0.8)
        )
        star.fill.solid()
        star.fill.fore_color.rgb = RGBColor(255, 215, 0)
        star.line.fill.background()
    
    # 底部文字
    footer = slide.shapes.add_textbox(
        Inches(2), Inches(5.8), Inches(9.333), Inches(1)
    )
    footer_tf = footer.text_frame
    footer_tf.text = "传承红色基因  描绘美好家园"
    footer_tf.paragraphs[0].font.size = Pt(28)
    footer_tf.paragraphs[0].font.color.rgb = RGBColor(255, 220, 180)
    footer_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 保存
    prs.save(output_path)
    print(f"PPT已保存至: {output_path}")


if __name__ == "__main__":
    output = os.path.expanduser("~/Desktop/毛泽东故居-小学生美术作品.pptx")
    create_ppt(output)
